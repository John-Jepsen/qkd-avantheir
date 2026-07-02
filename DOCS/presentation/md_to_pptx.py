"""Convert slides.md (Marp markdown) into an editable 16:9 PowerPoint deck.

Usage: python md_to_pptx.py [slides.md] [slides.pptx]

Handles the subset of markdown this deck uses: `---` slide separators,
#/## headings, bulleted and numbered lists, pipe tables, fenced code
blocks, and inline **bold** / `code` spans.
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DARK = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0x0E, 0x8A, 0x16)
CODE_BG = RGBColor(0xF3, 0xF4, 0xF6)
TABLE_HEADER = RGBColor(0x1F, 0x29, 0x37)
BODY_FONT = "Calibri"
CODE_FONT = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN


def parse_slides(md_text):
    """Split markdown into slides of typed blocks."""
    body = re.sub(r"\A---\n.*?\n---\n", "", md_text, flags=re.S)  # front matter
    slides = []
    for chunk in re.split(r"\n---\n", body):
        chunk = re.sub(r"<!--.*?-->", "", chunk, flags=re.S).strip()
        if not chunk:
            continue
        blocks, lines, i = [], chunk.split("\n"), 0
        while i < len(lines):
            line = lines[i]
            if line.startswith("```"):
                code = []
                i += 1
                while i < len(lines) and not lines[i].startswith("```"):
                    code.append(lines[i])
                    i += 1
                blocks.append(("code", "\n".join(code)))
            elif line.startswith("|"):
                rows = []
                while i < len(lines) and lines[i].startswith("|"):
                    cells = [c.strip() for c in lines[i].strip("|").split("|")]
                    if not all(re.fullmatch(r":?-+:?", c) for c in cells):
                        rows.append(cells)
                    i += 1
                blocks.append(("table", rows))
                continue
            elif m := re.match(r"^(#{1,2})\s+(.*)", line):
                blocks.append(("h1" if len(m.group(1)) == 1 else "h2", m.group(2)))
            elif m := re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)", line):
                indent, marker, text = m.groups()
                # merge hanging continuation lines into the bullet
                while i + 1 < len(lines) and re.match(r"^\s{2,}(?![-*\d])\S", lines[i + 1]):
                    i += 1
                    text += " " + lines[i].strip()
                blocks.append(("bullet", (len(indent) // 2, marker != "-" and marker != "*", text)))
            elif line.strip():
                text = line.strip()
                while i + 1 < len(lines) and lines[i + 1].strip() and not re.match(
                    r"^(#|[-*]\s|\d+\.\s|\||```)", lines[i + 1].lstrip()
                ):
                    i += 1
                    text += " " + lines[i].strip()
                blocks.append(("para", text))
            i += 1
        slides.append(blocks)
    return slides


def add_runs(paragraph, text, size, base_bold=False, color=DARK):
    """Split inline **bold** and `code` spans into styled runs."""
    for part in re.split(r"(\*\*.*?\*\*|`[^`]*`)", text):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("`") and part.endswith("`"):
            run.text = part[1:-1]
            run.font.name = CODE_FONT
            run.font.size = Pt(max(size - 2, 10))
        else:
            run.text = part.replace("*", "")
            run.font.bold = base_bold
        run.font.size = run.font.size or Pt(size)
        run.font.color.rgb = color
        if run.font.name is None:
            run.font.name = BODY_FONT


def build(md_path, out_path):
    slides = parse_slides(Path(md_path).read_text())
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for idx, blocks in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        y = MARGIN
        title_done = False
        for kind, data in blocks:
            if kind in ("h1", "h2") and not title_done:
                # first heading on the slide is the title
                h = Inches(1.5 if idx == 0 else 0.9)
                box = slide.shapes.add_textbox(MARGIN, y, CONTENT_W, h)
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                add_runs(p, data, 40 if idx == 0 else 30, base_bold=True)
                if idx == 0:
                    p.alignment = PP_ALIGN.CENTER
                    box.top = Inches(2.2)
                y += h + Inches(0.1)
                title_done = True
            elif kind in ("h1", "h2"):  # second heading mid-slide
                box = slide.shapes.add_textbox(MARGIN, y, CONTENT_W, Inches(0.6))
                p = box.text_frame.paragraphs[0]
                add_runs(p, data, 24, base_bold=True, color=ACCENT)
                y += Inches(0.7)
            elif kind == "bullet":
                level, numbered, text = data
                box = slide.shapes.add_textbox(
                    MARGIN + Inches(0.3 * level), y, CONTENT_W - Inches(0.3 * level), Inches(0.4)
                )
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                add_runs(p, ("• " if not numbered else "") + text, 17)
                box_h = Inches(0.34 * max(1, (len(text) // 95) + 1))
                box.height = box_h
                y += box_h + Emu(38100)
            elif kind == "para":
                box = slide.shapes.add_textbox(MARGIN, y, CONTENT_W, Inches(0.5))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                size = 20 if idx == 0 else 17
                add_runs(p, data, size)
                if idx == 0:
                    p.alignment = PP_ALIGN.CENTER
                    box.top = Inches(4.2)
                h = Inches(0.4 * max(1, (len(data) // 95) + 1))
                box.height = h
                y += h + Emu(38100)
            elif kind == "code":
                lines = data.split("\n")
                h = Inches(0.25 * len(lines) + 0.25)
                box = slide.shapes.add_textbox(MARGIN, y, CONTENT_W, h)
                box.fill.solid()
                box.fill.fore_color.rgb = CODE_BG
                tf = box.text_frame
                tf.word_wrap = False
                for j, line in enumerate(lines):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    run = p.add_run()
                    run.text = line if line else " "
                    run.font.name = CODE_FONT
                    run.font.size = Pt(13)
                    run.font.color.rgb = DARK
                y += h + Inches(0.15)
            elif kind == "table":
                rows, cols = len(data), len(data[0])
                h = Inches(0.42 * rows)
                shape = slide.shapes.add_table(rows, cols, MARGIN, y, CONTENT_W, h)
                table = shape.table
                for r, row in enumerate(data):
                    for c, cell_text in enumerate(row):
                        cell = table.cell(r, c)
                        cell.text = ""
                        p = cell.text_frame.paragraphs[0]
                        add_runs(
                            p,
                            cell_text,
                            14,
                            base_bold=(r == 0),
                            color=RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else DARK,
                        )
                y += h + Inches(0.2)

    prs.save(out_path)
    print(f"wrote {out_path} ({len(slides)} slides)")


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else "slides.md"
    dst = sys.argv[2] if len(sys.argv) > 2 else "slides.pptx"
    build(src, dst)
